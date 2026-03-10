#!/usr/bin/env python3
"""
Build clean slide-master background images from a template PPTX.

Workflow:
1) Copy input template to a temp PPTX.
2) Remove placeholder prompt text (e.g., "Click to edit master text styles")
   on selected slides.
3) Convert temp PPTX to PDF using soffice.
4) Render selected PDF pages to PNG backgrounds.

Default output names match this repo's styleguide:
- bg-title-master.png
- bg-content-master.png
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Optional, Tuple


PLACEHOLDER_PATTERNS = [
    re.compile(r"click\s+to\s+edit\s+master\s+text\s+styles", re.IGNORECASE),
    re.compile(r"click\s+to\s+add\s+title", re.IGNORECASE),
    re.compile(r"click\s+to\s+add\s+text", re.IGNORECASE),
    re.compile(r"click\s+to\s+add\s+sub.?title", re.IGNORECASE),
    re.compile(r"click\s+to\s+add\s+.*", re.IGNORECASE),
]


def _is_placeholder_prompt(text: str) -> bool:
    t = re.sub(r"\s+", " ", (text or "")).strip()
    if not t:
        return False
    return any(p.search(t) for p in PLACEHOLDER_PATTERNS)


def _clear_placeholder_prompts(pptx_path: Path, slide_numbers_1based: list[int]) -> int:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(f"python-pptx is required: {exc}") from exc

    prs = Presentation(str(pptx_path))
    cleared = 0

    for slide_no in slide_numbers_1based:
        idx = slide_no - 1
        if idx < 0 or idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = tf.text or ""
            if not _is_placeholder_prompt(text):
                continue
            tf.clear()
            cleared += 1

    prs.save(str(pptx_path))
    return cleared


def _find_soffice() -> str:
    soffice = shutil.which("soffice")
    if soffice:
        return soffice
    fallback = Path.home() / ".local" / "bin" / "soffice"
    if fallback.exists():
        return str(fallback)
    raise RuntimeError("soffice not found. Install LibreOffice and ensure 'soffice' is on PATH.")


def _convert_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    soffice = _find_soffice()
    subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(pptx_path),
        ],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    pdf = out_dir / f"{pptx_path.stem}.pdf"
    if not pdf.exists():
        raise RuntimeError(f"Expected PDF not found: {pdf}")
    return pdf


def _render_pdf_page_to_png(pdf_path: Path, page_1based: int, out_png: Path, zoom: float = 2.0) -> None:
    try:
        import fitz  # PyMuPDF
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(f"PyMuPDF is required: {exc}") from exc

    doc = fitz.open(str(pdf_path))
    idx = page_1based - 1
    if idx < 0 or idx >= len(doc):
        raise RuntimeError(f"PDF page {page_1based} out of range (1..{len(doc)})")
    page = doc[idx]
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)

    from PIL import Image
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    img = _remove_bottom_right_slide_number(page=page, pil_img=img, zoom=zoom)

    out_png.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(out_png))


def _detect_bottom_right_number_bbox(page) -> Optional[Tuple[float, float, float, float]]:
    """Find likely slide-number text bbox in the lower-right corner in page units."""
    r = page.rect
    # Search a conservative corner region to avoid wiping real content.
    clip = type(r)(r.x1 - (r.width * 0.14), r.y1 - (r.height * 0.14), r.x1, r.y1)
    words = page.get_text("words", clip=clip) or []
    if not words:
        return None

    candidates = []
    for w in words:
        if len(w) < 5:
            continue
        x0, y0, x1, y1, token = w[:5]
        t = str(token or "").strip()
        if not re.fullmatch(r"\d{1,3}", t):
            continue
        # Prefer smallest/lowest-right number token in the corner.
        area = max(0.0, float(x1) - float(x0)) * max(0.0, float(y1) - float(y0))
        candidates.append((float(x0), float(y0), float(x1), float(y1), area))

    if not candidates:
        return None

    candidates.sort(key=lambda c: (c[1] + c[3], c[0] + c[2], c[4]))
    x0, y0, x1, y1, _ = candidates[-1]
    return (x0, y0, x1, y1)


def _sample_fill_color(pil_img, px0: int, py0: int, px1: int, py1: int):
    """Sample a nearby area color to blend patch with background."""
    from PIL import ImageStat

    w, h = pil_img.size
    # Sample just left of the number box first.
    sx0 = max(0, px0 - (px1 - px0) - 10)
    sx1 = max(0, px0 - 4)
    sy0 = max(0, py0 - 2)
    sy1 = min(h, py1 + 2)
    if sx1 > sx0 and sy1 > sy0:
        stat = ImageStat.Stat(pil_img.crop((sx0, sy0, sx1, sy1)))
        return tuple(int(v) for v in stat.mean[:3])

    # Fallback sample from above if left sample is invalid.
    sx0 = max(0, px0 - 2)
    sx1 = min(w, px1 + 2)
    sy0 = max(0, py0 - (py1 - py0) - 10)
    sy1 = max(0, py0 - 4)
    if sx1 > sx0 and sy1 > sy0:
        stat = ImageStat.Stat(pil_img.crop((sx0, sy0, sx1, sy1)))
        return tuple(int(v) for v in stat.mean[:3])

    return (255, 255, 255)


def _remove_bottom_right_slide_number(page, pil_img, zoom: float):
    """Mask detected slide number in lower-right corner from rendered image."""
    bbox = _detect_bottom_right_number_bbox(page)
    if not bbox:
        return pil_img

    from PIL import ImageDraw

    x0, y0, x1, y1 = bbox
    px0 = int(max(0, x0 * zoom))
    py0 = int(max(0, y0 * zoom))
    px1 = int(max(0, x1 * zoom))
    py1 = int(max(0, y1 * zoom))

    # Add a small pad to fully cover anti-aliased glyph edges.
    pad_x = 6
    pad_y = 4
    px0 = max(0, px0 - pad_x)
    py0 = max(0, py0 - pad_y)
    w, h = pil_img.size
    px1 = min(int(w), px1 + pad_x)
    py1 = min(int(h), py1 + pad_y)
    if px1 <= px0 or py1 <= py0:
        return pil_img

    fill = _sample_fill_color(pil_img, px0, py0, px1, py1)
    draw = ImageDraw.Draw(pil_img)
    draw.rectangle((px0, py0, px1, py1), fill=fill)
    return pil_img


def build_backgrounds(
    input_pptx: Path,
    output_dir: Path,
    title_slide: int,
    content_slide: int,
    title_name: str,
    content_name: str,
    third_slide: Optional[int] = None,
    third_name: Optional[str] = None,
) -> dict:
    output_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="template-bg-") as td:
        tmp_dir = Path(td)
        tmp_pptx = tmp_dir / input_pptx.name
        shutil.copy2(input_pptx, tmp_pptx)

        slide_targets = [title_slide, content_slide]
        if third_slide is not None:
            slide_targets.append(third_slide)
        cleared = _clear_placeholder_prompts(tmp_pptx, slide_targets)
        pdf_path = _convert_pptx_to_pdf(tmp_pptx, tmp_dir)

        title_png = output_dir / title_name
        content_png = output_dir / content_name
        _render_pdf_page_to_png(pdf_path, title_slide, title_png)
        _render_pdf_page_to_png(pdf_path, content_slide, content_png)
        third_png = None
        if third_slide is not None and third_name:
            third_png = output_dir / third_name
            _render_pdf_page_to_png(pdf_path, third_slide, third_png)

    result = {
        "input": str(input_pptx),
        "outputDir": str(output_dir),
        "titleSlide": title_slide,
        "contentSlide": content_slide,
        "titleBackground": str(title_png),
        "contentBackground": str(content_png),
        "placeholderPromptsCleared": cleared,
    }
    if third_slide is not None and third_name and third_png is not None:
        result["thirdSlide"] = third_slide
        result["thirdBackground"] = str(third_png)
    return result


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build clean PPTX background images from template slides.")
    parser.add_argument("--input", required=True, help="Path to source template .pptx")
    parser.add_argument(
        "--output-dir",
        default="assets/freeport-template",
        help="Output directory for background PNG assets",
    )
    parser.add_argument("--title-slide", type=int, default=1, help="1-based slide number for title master capture")
    parser.add_argument("--content-slide", type=int, default=2, help="1-based slide number for content master capture")
    parser.add_argument("--title-name", default="bg-title-master.png", help="Output filename for title background")
    parser.add_argument("--content-name", default="bg-content-master.png", help="Output filename for content background")
    parser.add_argument("--third-slide", type=int, default=None, help="Optional 1-based slide number for third master capture")
    parser.add_argument("--third-name", default="bg-reconciliation-master.png", help="Output filename for third background")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    input_pptx = Path(args.input).expanduser().resolve()
    output_dir = Path(args.output_dir).expanduser().resolve()

    if not input_pptx.exists():
        print(f"Input file not found: {input_pptx}", file=sys.stderr)
        return 2
    if input_pptx.suffix.lower() != ".pptx":
        print("--input must be a .pptx file", file=sys.stderr)
        return 2

    try:
        result = build_backgrounds(
            input_pptx=input_pptx,
            output_dir=output_dir,
            title_slide=args.title_slide,
            content_slide=args.content_slide,
            title_name=args.title_name,
            content_name=args.content_name,
            third_slide=args.third_slide,
            third_name=args.third_name,
        )
    except Exception as exc:
        print(f"Template background build failed: {exc}", file=sys.stderr)
        return 1

    print(json.dumps(result, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
